from pptx import Presentation
import os
from pptx.util import Pt
from params import *
from functions import obtenerpath, wrap_label, procesar_dataframe_simbolos, procesar_dataframe_circulos, obtenerpath2
#, get_color_for_value
import numpy as np
#import seaborn as sns
#import matplotlib.pyplot as plt
import random
import pandas as pd

def create_memorando(tipo,negocio,codigo,df):

    

    if negocio == 'SEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "07. Sprint Planning.pptx")
    elif negocio == 'PRIMA':
        ppt_template_path = os.path.join(PLANTILLAS_PRI_PATH, "07. Sprint Planning.pptx")
    
    elif negocio == 'SALUD':
        ppt_template_path = os.path.join(PLANTILLAS_EPS_PATH, "07. Sprint Planning.pptx")
    
    elif negocio == 'CREDISEGURO':
        ppt_template_path = os.path.join(PLANTILLAS_CRE_PATH, "07. Sprint Planning.pptx")
    else:
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "07. Sprint Planning.pptx")
    


    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    for slide_num, slide_data in enumerate(memorando_list, start=1):
        if slide_num > len(ppt.slides):
            print(f"No hay suficiente cantidad de slides en la presentación para el diccionario número {slide_num}")
            break
        slide = ppt.slides[slide_num - 1]
        print(f"Diapositiva número: {slide_num}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, column_name in slide_data.items():
                            if run.text == key:
                                run.text = str(df[column_name].iloc[0])

    try:
        final_path = obtenerpath(tipo,negocio,codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")




def create_observaciones_ppt(tipo, negocio, codigo, df):
    num_slides = len(df) 
    
    if negocio == 'SEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    elif negocio == 'PRIMA':
        ppt_template_path = os.path.join(PLANTILLAS_PRI_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    
    elif negocio == 'SALUD':
        ppt_template_path = os.path.join(PLANTILLAS_EPS_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    
    elif negocio == 'CREDISEGURO':
        ppt_template_path = os.path.join(PLANTILLAS_CRE_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    else:
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, f"plantilla_observaciones_{num_slides}.pptx")

    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

   
    first_slide = ppt.slides[0]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in primera_slide_dict.items():
                        if key in run.text:
                            run.text = run.text.replace(key, str(value))

    
    font_size_observaciones = Pt(14)

    
    for index, row in enumerate(df.itertuples(index=False)):
        slide = ppt.slides[index + 1]  
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, column_name in observaciones_list[0].items():
                            if key in run.text:
                                if key == 'Unidad_auditada_n':
                                    run.text = run.text.replace(key, str(getattr(row, column_name)))
                                else:
                                    run.text = run.text.replace(key, str(getattr(row, column_name)))
                                    run.font.size = font_size_observaciones

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    for key, column_name in observaciones_list[0].items():
                        if key in cell.text:
                            cell.text = cell.text.replace(key, str(getattr(row, column_name)))
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = font_size_observaciones

    try:
        final_path = obtenerpath(tipo, negocio, codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")


def create_informe_ppt(tipo, negocio, codigo, df_datos, df_calificativo_total, df_calificativo_unidad_responsable, df_cantidad_controles, df_observaciones_informe):

    if negocio == 'SEGUROS':
        if tipo == 'Informe Final':
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Informe Final.pptx")
        else:
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")

    elif negocio == 'PRIMA':
        if tipo == 'Informe Final':
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Informe Final.pptx")
        else:
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")
    elif negocio == 'SALUD':
        if tipo == 'Informe Final':
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Informe Final.pptx")
        else:
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")
    elif negocio == 'CREDISEGURO':
        if tipo == 'Informe Final':
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Informe Final.pptx")
        else:
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")
    else:
        if tipo == 'Informe Final':
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Informe Final.pptx")
        else:
            ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")

    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    listas_diccionarios = [
        (calificativo_total_list, df_calificativo_total),
        (calificativo_unidad_responsable_list, df_calificativo_unidad_responsable),
        (cantidad_controles_list, df_cantidad_controles),
        (observaciones_informe_list, df_observaciones_informe),
        (memorando_list, df_datos)
    ]

    font_size_first_slide = Pt(40)
    font_size_other_slides = Pt(18)

    for lista_dicc, df in listas_diccionarios:
        for slide_num, slide_data in enumerate(lista_dicc):
            if slide_num >= len(ppt.slides):
                break
            slide = ppt.slides[slide_num]
            font_size = font_size_first_slide if slide_num == 0 else font_size_other_slides

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key, column_name in slide_data.items():
                                if lista_dicc in [calificativo_unidad_responsable_list, observaciones_informe_list]:
                                    max_records = 3 if lista_dicc == observaciones_informe_list else 2
                                    for i in range(min(len(df), max_records)):
                                        if not df.empty and column_name in df.columns:
                                            new_key = f"{key}_{i + 1}"
                                            if new_key in run.text:
                                                run.text = run.text.replace(new_key, str(df[column_name].iloc[i]))
                                                run.font.size = font_size
                                else:
                                    if not df.empty and column_name in df.columns:
                                        if key in run.text:
                                            run.text = run.text.replace(key, str(df[column_name].iloc[0]))
                                            run.font.size = font_size

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for cell in table.iter_cells():
                        for key, column_name in slide_data.items():
                            if lista_dicc in [calificativo_unidad_responsable_list, observaciones_informe_list]:
                                max_records = 3 if lista_dicc == observaciones_informe_list else 2
                                for i in range(min(len(df), max_records)):
                                    if not df.empty and column_name in df.columns:
                                        new_key = f"{key}_{i + 1}"
                                        if new_key in cell.text:
                                            cell.text = cell.text.replace(new_key, str(df[column_name].iloc[i]))
                                            for paragraph in cell.text_frame.paragraphs:
                                                for run in paragraph.runs:
                                                    run.font.size = font_size
                            else:
                                if not df.empty and column_name in df.columns:
                                    if key in cell.text:
                                        cell.text = cell.text.replace(key, str(df[column_name].iloc[0]))
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = font_size

    try:
        final_path = obtenerpath(tipo, negocio, codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")




'''
def create_mapa_calor(tipo,codigo,df_simbolos,df_circulos):
    
    data = np.array([
        [2, 2, 3, 3, 3],
        [1, 2, 3, 3, 3],
        [1, 2, 2, 3, 3],
        [1, 1, 2, 2, 2],
        [1, 1, 1, 1, 2]
    ])

    
    x_labels = ["Bajo (1)", "Moderado (2)", "Relevante (3)", "Alto (4)", "Crítico (5)"]
    y_labels = [
        "Efectivo (1)",
        "Efectivo No Formalizado (2)",
        "Inefectivo de Prueba (3)",
        "Inefectivo de Diseño (4)",
        "Control Inexistente (5)",
    ]

    y_labels_wrapped = [wrap_label(label, max_width=17) for label in y_labels]

  
    cmap = sns.color_palette(["green", "yellow", "red"])

    
    fig, ax = plt.subplots(figsize=(8, 6)) 

    sns.heatmap(
        data,
        annot=False,
        cmap=cmap,
        linewidths=0.5,
        linecolor="black",
        xticklabels=x_labels,
        yticklabels=False,
        cbar=False,
        square=False
    )

        
    for i, label in enumerate(y_labels_wrapped):
        ax.text(
            -0.1,
            len(data) - i - 0.5,
            label,
            rotation=90,
            va="center",
            ha="center",
            fontsize=6.5
        )

    
    ax.text(
        -0.3,
        len(data) / 2,
        "Administración del Riesgo",
        rotation=90,
        va="center",
        ha="center",
        fontsize=8.5,
        fontweight="bold"
    )

    
    cantidades_circulos = procesar_dataframe_circulos(df_circulos)

    
    cantidades_simbolos = procesar_dataframe_simbolos(df_simbolos)

    
    simbolos_tipo = {
        "Riesgo de Crédito": "#",
        "Riesgo Operacional": "▲",
        "Riesgo de Mercado": "■",
        "Riesgo Legal": "$",
        "Riesgo Estratégico": "x",
        "Riesgo de Producción/Suscripción": "◆"
    }
    total_controles = sum(cantidades_circulos.values())
    numeros_controles = random.sample(range(1, total_controles + 1), total_controles)

   
    circle_step = 0.32 
    circle_size = 11
    separacion = 0.1

    
    for (y, x), cantidad in cantidades_circulos.items():
        coords = []
        max_circles = 9  

        for i in range(min(cantidad, max_circles)):
            dx = (i % 3 - 1) * circle_step
            dy = (i // 3 - 1) * circle_step
            coords.append((x + 0.5 + dx, y + 0.5 + dy))

        for i in range(max_circles, cantidad):


            coords.append((x + 0.2+separacion, y + 0.2+separacion))
            separacion += 0.15

        for coord in coords:
            if numeros_controles:
                numero = numeros_controles.pop(0)
                ax.text(
                    coord[0], coord[1],
                    f"{numero}",
                    ha="center",
                    va="center",
                    fontsize=circle_size,
                    color="white",
                    bbox=dict(boxstyle="circle", facecolor="blue", edgecolor="blue", pad=0.4)
                )

    ocupados = {}

    
    agrupados_por_cuadrante = {}
    for ((y, x), riesgo), cantidad in cantidades_simbolos.items():
        if (y, x) not in agrupados_por_cuadrante:
            agrupados_por_cuadrante[(y, x)] = []
        agrupados_por_cuadrante[(y, x)].append((riesgo, cantidad))

    for (y, x), categorias in agrupados_por_cuadrante.items():
        if (y, x) not in ocupados:
            ocupados[(y, x)] = []

        
        base_x, base_y = x + 0.5, y + 0.5
        step = 0.2  
        simbolos_por_cuadrante = len(categorias)

       
        for idx, (riesgo, cantidad) in enumerate(categorias):
            simbolo = simbolos_tipo[riesgo]

            for i in range(cantidad):
                
                dx = (i % 3 - 1) * step + (idx % 2) * 0.1 
                dy = (i // 3 - 1) * step + (idx // 2) * 0.1
                coord = (base_x + dx, base_y + dy)

                if coord not in ocupados[(y, x)]:
                    ocupados[(y, x)].append(coord)
                    ax.text(
                        coord[0], coord[1],
                        simbolo,
                        ha="center",
                        va="center",
                        fontsize=14,
                        color="white", 
                    )
    
    plt.xlabel("Riesgo Inherente", fontsize=8.5, fontweight="bold", labelpad=4)
    plt.xticks(rotation=0, fontsize=7)

    
    plt.subplots_adjust(left=0.02)
    plt.tight_layout()

    

    try:
        final_path = obtenerpath2(codigo)
        codigo_final = codigo.replace(' ','')
        nombre = f"Mapa_de_calor_{codigo_final}.png"
        path_final= os.path.join(final_path,nombre)
        plt.savefig(path_final)
        print(f"imagen guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")

    
    plt.close()


'''


"""
   
  def create_evolutivo(id_evaluacion, file_path='./Input_PlanAnual.xlsx'):
    
    df = pd.read_excel(file_path)
    
    
    filtered_df = df[(df['ID'] == id_evaluacion)]
    if filtered_df.empty:
        print("No se encontró la evaluación y negocio especificados:", id_evaluacion)
        return
    
   
    row = filtered_df.iloc[0]
    
    
    year_columns = ['2019', '2020', '2021', '2022', '2023', '2024']
    years_available = []
    effective_values = [] 
    display_values = []   
    
    for col in year_columns:
        if col in row:
            val = row[col]
            if pd.notnull(val):
                
                if isinstance(val, str) and val.strip().upper() == "S/C":
                    years_available.append(int(col))
                    effective_values.append(1.0)
                    display_values.append("S/C")
                else:
                    try:
                        numeric_val = float(val)
                        years_available.append(int(col))
                        effective_values.append(numeric_val)
                        display_values.append(f'{numeric_val:.2f}')
                    except ValueError:
                       
                        pass
    
    if not effective_values:
        print("No hay calificativos disponibles para la evaluación:", evaluation_name, "y negocio:", negocio)
        return
    
 
    if len(effective_values) >= 3:
        sel_years = years_available[-3:]
        sel_effective = effective_values[-3:]
        sel_display = display_values[-3:]
    elif len(effective_values) == 2:
        sel_years = years_available[-2:]
        sel_effective = effective_values[-2:]
        sel_display = display_values[-2:]
    else:
        sel_years = years_available[-1:]
        sel_effective = effective_values[-1:]
        sel_display = display_values[-1:]
    
    
    plt.figure(figsize=(8, 5))
    ax = plt.gca()
    
    plt.plot(sel_years, sel_effective, color='gray', linewidth=2, zorder=1)
    
    for year, eff_val, disp_val in zip(sel_years, sel_effective, sel_display):
        color = get_color_for_value(eff_val)
        plt.scatter(year, eff_val, color=color, s=80, zorder=2)
        plt.text(year, eff_val + 0.02, disp_val, ha='center', va='bottom', fontsize=9, color=color, zorder=3)
    
   
    plt.axhspan(0, 0.55, color='red', alpha=0.1)
    plt.axhspan(0.55, 0.70, color='orange', alpha=0.1)
    plt.axhspan(0.70, 0.85, color='green', alpha=0.1)
    plt.axhspan(0.85, 1.0, color='lightgreen', alpha=0.2)
    
    plt.ylim(0, 1.05)
    plt.xlim(min(sel_years) - 0.5, max(sel_years) + 0.5)
    plt.xticks(sel_years)
   
    plt.xlabel('Año')
    plt.ylabel('Calificativo (decimal)')
    
    plt.tight_layout()
    
   
    filename = f'evolucion_{evaluation_name}_{negocio}.png'
    plt.savefig(filename, dpi=300, bbox_inches='tight')
    
    print(f'Gráfico guardado como {filename}')
 
   
   
   
   
""" 
